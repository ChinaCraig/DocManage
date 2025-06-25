"""
æ–‡æ¡£ç”ŸæˆæœåŠ¡ - åŸºäºç°æœ‰æ–‡ä»¶/æ–‡ä»¶å¤¹ç”Ÿæˆæ–°æ–‡æ¡£
ç‰ˆæœ¬: v1.0
åŠŸèƒ½: æ”¯æŒåŸºäºæ–‡ä»¶å¤¹æˆ–æ–‡ä»¶å†…å®¹ç”Ÿæˆtxtæ–‡æ¡£ï¼Œä½¿ç”¨æ¨¡å—åŒ–è®¾è®¡ä¾¿äºæ‰©å±•
"""

import os
import re
import json
import logging
import time
from typing import Dict, Any, List, Optional, Union
from datetime import datetime
from sqlalchemy import and_, or_

from app.models.document_models import DocumentNode, DocumentContent, db
from app.services.llm import LLMService
from app.services.image_recognition_service import ImageRecognitionService

logger = logging.getLogger(__name__)

class DocumentGenerationService:
    """æ–‡æ¡£ç”ŸæˆæœåŠ¡"""
    
    # æ”¯æŒçš„è¾“å‡ºæ ¼å¼
    SUPPORTED_FORMATS = ['txt', 'pdf', 'xlsx', 'docx']  # æ”¯æŒæ–‡æœ¬ã€PDFã€Excelã€Wordæ ¼å¼
    
    # æ”¯æŒçš„æ–‡æ¡£ç±»å‹
    DOCUMENT_TYPES = {
        'summary': 'æ€»ç»“æ±‡æ€»',
        'report': 'åˆ†ææŠ¥å‘Š', 
        'analysis': 'æ·±åº¦åˆ†æ',
        'documentation': 'è¯´æ˜æ–‡æ¡£',
        'other': 'å…¶ä»–ç±»å‹'
    }
    
    @staticmethod
    def generate_document(source_path: str, output_format: str = 'txt', 
                         document_type: str = 'summary', 
                         query_text: str = '',
                         llm_model: str = None) -> Dict[str, Any]:
        """
        ç”Ÿæˆæ–‡æ¡£ä¸»å…¥å£
        
        Args:
            source_path: æºæ–‡ä»¶æˆ–æ–‡ä»¶å¤¹è·¯å¾„/åç§°
            output_format: è¾“å‡ºæ ¼å¼ (å½“å‰åªæ”¯æŒtxt)
            document_type: æ–‡æ¡£ç±»å‹
            query_text: åŸå§‹æŸ¥è¯¢æ–‡æœ¬
            llm_model: LLMæ¨¡å‹
            
        Returns:
            Dict: ç”Ÿæˆç»“æœ
        """
        try:
            # 0. æ ‡å‡†åŒ–è¾“å‡ºæ ¼å¼ï¼ˆå…¼å®¹æ„å›¾è¯†åˆ«æœåŠ¡çš„æ ¼å¼åç§°ï¼‰
            format_mapping = {
                'excel': 'xlsx',  # Excelæ ¼å¼ç»Ÿä¸€ä¸ºxlsx
                'doc': 'docx',    # Wordæ ¼å¼ç»Ÿä¸€ä¸ºdocx
                'word': 'docx',   # Wordæ ¼å¼çš„å¦ä¸€ç§è¡¨è¾¾
                'ppt': 'pptx',    # PowerPointæ ¼å¼ï¼ˆæš‚æœªå®ç°ï¼‰
                'text': 'txt',    # æ–‡æœ¬æ ¼å¼çš„å¦ä¸€ç§è¡¨è¾¾
            }
            
            # åº”ç”¨æ ¼å¼æ˜ å°„
            original_format = output_format
            output_format = format_mapping.get(output_format.lower(), output_format.lower())
            
            if original_format != output_format:
                logger.info(f"æ ¼å¼æ ‡å‡†åŒ–: {original_format} â†’ {output_format}")
            
            # 1. éªŒè¯è¾“å‡ºæ ¼å¼
            if output_format not in DocumentGenerationService.SUPPORTED_FORMATS:
                return {
                    'success': False,
                    'error': f'æš‚ä¸æ”¯æŒ {original_format} æ ¼å¼ï¼Œå½“å‰åªæ”¯æŒ: {", ".join(DocumentGenerationService.SUPPORTED_FORMATS)}',
                    'query': query_text
                }
            
            # 2. æŸ¥æ‰¾æºæ–‡ä»¶æˆ–æ–‡ä»¶å¤¹ï¼ˆæ”¯æŒæ¨¡ç³ŠæŸ¥è¯¢ï¼‰
            source_node = DocumentGenerationService._find_source_node(source_path)
            if not source_node:
                # å¦‚æœæ‰¾ä¸åˆ°æºæ–‡ä»¶ï¼Œå°è¯•æ™ºèƒ½æ£€ç´¢
                logger.info(f"æœªæ‰¾åˆ°æºæ–‡ä»¶æˆ–æ–‡ä»¶å¤¹: '{source_path}'ï¼Œå°è¯•æ™ºèƒ½æ£€ç´¢...")
                
                try:
                    # ä½¿ç”¨æºè·¯å¾„ä½œä¸ºæœç´¢å…³é”®è¯è¿›è¡Œå‘é‡æ£€ç´¢
                    from app.services.vectorization.vector_service_adapter import VectorServiceAdapter
                    
                    vector_adapter = VectorServiceAdapter()
                    search_results = vector_adapter.search(
                        query_text=source_path,  # ç›´æ¥ä½¿ç”¨æºè·¯å¾„ä½œä¸ºæœç´¢å…³é”®è¯
                        top_k=10,
                        min_score=0.3
                    )
                    
                    if search_results and len(search_results) > 0:
                        # ä½¿ç”¨æœç´¢ç»“æœç”Ÿæˆæ–‡æ¡£
                        logger.info(f"æ™ºèƒ½æ£€ç´¢æ‰¾åˆ° {len(search_results)} ä¸ªç›¸å…³ç»“æœï¼Œä½¿ç”¨æœç´¢ç»“æœç”Ÿæˆæ–‡æ¡£")
                        
                        return DocumentGenerationService.generate_from_search_results(
                            search_results=search_results,
                            output_format=output_format,
                            document_type=document_type,
                            query_text=query_text,
                            llm_model=llm_model,
                            search_keywords=source_path
                        )
                    else:
                        return {
                            'success': False,
                            'error': f'æœªæ‰¾åˆ°ä¸"{source_path}"ç›¸å…³çš„æ–‡æ¡£æˆ–æ–‡ä»¶ï¼Œè¯·å°è¯•ä½¿ç”¨æ›´å…·ä½“çš„æ–‡ä»¶åæˆ–æ–‡ä»¶å¤¹å',
                            'query': query_text,
                            'source_path': source_path,
                            'suggestion': 'ğŸ’¡ å»ºè®®ï¼šå¯ä»¥å°è¯•ä½¿ç”¨å®Œæ•´çš„æ–‡ä»¶åæˆ–æ–‡ä»¶å¤¹åï¼Œä¾‹å¦‚ï¼š"åŸºäºé”€å”®æ•°æ®æ–‡ä»¶å¤¹ç”ŸæˆæŠ¥å‘Š"'
                        }
                        
                except Exception as e:
                    logger.error(f"æ™ºèƒ½æ£€ç´¢å¤±è´¥: {e}")
                    return {
                        'success': False,
                        'error': f'æœªæ‰¾åˆ°æºæ–‡ä»¶æˆ–æ–‡ä»¶å¤¹: "{source_path}"ï¼Œæ™ºèƒ½æ£€ç´¢ä¹Ÿæœªæ‰¾åˆ°ç›¸å…³å†…å®¹',
                        'query': query_text,
                        'source_path': source_path
                    }
            
            # 3. æ ¹æ®æºç±»å‹å¤„ç†
            if source_node.type == 'folder':
                return DocumentGenerationService._generate_from_folder(
                    source_node, output_format, document_type, query_text, llm_model
                )
            else:
                return DocumentGenerationService._generate_from_file(
                    source_node, output_format, document_type, query_text, llm_model
                )
                
        except Exception as e:
            logger.error(f"æ–‡æ¡£ç”Ÿæˆå¤±è´¥: {e}")
            return {
                'success': False,
                'error': f'æ–‡æ¡£ç”Ÿæˆè¿‡ç¨‹ä¸­å‡ºç°é”™è¯¯: {str(e)}',
                'query': query_text,
                'source_path': source_path
            }
    
    @staticmethod
    def generate_from_search_results(search_results: List[Dict], output_format: str = 'txt',
                                   document_type: str = 'summary', query_text: str = '',
                                   llm_model: str = None, search_keywords: str = '') -> Dict[str, Any]:
        """
        åŸºäºæœç´¢ç»“æœç”Ÿæˆæ–‡æ¡£
        
        Args:
            search_results: èšåˆåçš„æ–‡ä»¶çº§æœç´¢ç»“æœ
            output_format: è¾“å‡ºæ ¼å¼
            document_type: æ–‡æ¡£ç±»å‹
            query_text: åŸå§‹æŸ¥è¯¢æ–‡æœ¬
            llm_model: LLMæ¨¡å‹
            search_keywords: æœç´¢å…³é”®è¯ï¼ˆç”¨äºæ–‡ä»¶å‘½åï¼‰
            
        Returns:
            Dict: ç”Ÿæˆç»“æœ
        """
        try:
            logger.info(f"å¼€å§‹åŸºäºæœç´¢ç»“æœç”Ÿæˆ {document_type} æ–‡æ¡£ï¼Œæ‰¾åˆ° {len(search_results)} ä¸ªç›¸å…³æ–‡ä»¶")
            
            # 0. æ ‡å‡†åŒ–è¾“å‡ºæ ¼å¼
            format_mapping = {
                'excel': 'xlsx', 'doc': 'docx', 'word': 'docx',
                'ppt': 'pptx', 'text': 'txt'
            }
            output_format = format_mapping.get(output_format.lower(), output_format.lower())
            
            # 1. éªŒè¯è¾“å‡ºæ ¼å¼
            if output_format not in DocumentGenerationService.SUPPORTED_FORMATS:
                return {
                    'success': False,
                    'error': f'æš‚ä¸æ”¯æŒ {output_format} æ ¼å¼ï¼Œå½“å‰åªæ”¯æŒ: {", ".join(DocumentGenerationService.SUPPORTED_FORMATS)}',
                    'query': query_text
                }
            
            # 2. è½¬æ¢æœç´¢ç»“æœä¸ºæ–‡ä»¶å†…å®¹æ ¼å¼
            file_contents = []
            processed_files = []
            
            for file_result in search_results[:10]:  # é™åˆ¶æœ€å¤š10ä¸ªæ–‡ä»¶é¿å…tokenè¶…é™
                try:
                    document_info = file_result.get('document', {})
                    document_id = document_info.get('id')
                    file_name = document_info.get('name', 'æœªçŸ¥æ–‡ä»¶')
                    
                    if document_id and document_id not in processed_files:
                        # åˆå¹¶è¯¥æ–‡ä»¶çš„æ‰€æœ‰åŒ¹é…ç‰‡æ®µ
                        chunks = file_result.get('chunks', [])
                        combined_content = '\n'.join([chunk.get('text', '') for chunk in chunks[:5]])  # æ¯ä¸ªæ–‡ä»¶æœ€å¤š5ä¸ªç‰‡æ®µ
                        
                        if combined_content:
                            file_contents.append({
                                'file_name': file_name,
                                'file_type': document_info.get('file_type', 'æœªçŸ¥'),
                                'document_id': document_id,
                                'content': combined_content[:2000],  # é™åˆ¶å†…å®¹é•¿åº¦
                                'relevance_score': file_result.get('average_score', 0)
                            })
                            processed_files.append(document_id)
                            
                except Exception as e:
                    logger.warning(f"å¤„ç†æœç´¢ç»“æœæ–‡ä»¶å¤±è´¥: {e}")
                    continue
            
            if not file_contents:
                return {
                    'success': False,
                    'error': 'æœç´¢ç»“æœä¸­æ²¡æœ‰å¯ç”¨çš„æ–‡ä»¶å†…å®¹',
                    'query': query_text,
                    'search_results_count': len(search_results)
                }
            
            # 3. ä½¿ç”¨LLMç”Ÿæˆæ–‡æ¡£
            generated_content = DocumentGenerationService._generate_content_with_llm(
                source_type='æœç´¢ç»“æœ',
                source_name=f"å…³é”®è¯'{search_keywords}'çš„æœç´¢ç»“æœ",
                source_contents=file_contents,
                document_type=document_type,
                query_text=query_text,
                llm_model=llm_model
            )
            
            # 4. ä¿å­˜ç”Ÿæˆçš„æ–‡æ¡£
            source_name_for_file = search_keywords if search_keywords else 'æœç´¢ç»“æœ'
            saved_file = DocumentGenerationService._save_generated_document(
                content=generated_content,
                source_name=source_name_for_file,
                document_type=document_type,
                output_format=output_format,
                parent_folder=None,  # æœç´¢ç»“æœç”Ÿæˆçš„æ–‡æ¡£æ”¾åœ¨æ ¹ç›®å½•
                query_text=query_text
            )
            
            return {
                'success': True,
                'query': query_text,
                'source_type': 'æœç´¢ç»“æœ',
                'source_info': {
                    'name': f"åŸºäºå…³é”®è¯'{search_keywords}'çš„æœç´¢ç»“æœ",
                    'search_keywords': search_keywords,
                    'total_files_found': len(search_results),
                    'files_used': len(file_contents)
                },
                'processed_files_count': len(file_contents),
                'total_files_count': len(search_results),
                'document_type': document_type,
                'output_format': output_format,
                'generated_content': generated_content,
                'saved_file': saved_file.to_dict() if saved_file else None,
                'generation_method': 'llm_search_results_analysis',
                'timestamp': datetime.now().isoformat()
            }
            
        except Exception as e:
            logger.error(f"åŸºäºæœç´¢ç»“æœç”Ÿæˆæ–‡æ¡£å¤±è´¥: {e}")
            return {
                'success': False,
                'error': f'åŸºäºæœç´¢ç»“æœç”Ÿæˆæ–‡æ¡£å¤±è´¥: {str(e)}',
                'query': query_text
            }
    
    @staticmethod
    def _find_source_node(source_path: str) -> Optional[DocumentNode]:
        """æŸ¥æ‰¾æºæ–‡ä»¶æˆ–æ–‡ä»¶å¤¹èŠ‚ç‚¹ - æ”¯æŒå¤šå…³é”®è¯æ¨¡ç³ŠæŸ¥è¯¢"""
        try:
            # æ¸…ç†è·¯å¾„åç§°
            source_name = source_path.strip()
            
            # ç§»é™¤å¸¸è§çš„ä¿®é¥°è¯
            source_name = re.sub(r'^(?:è¿™ä¸ª|é‚£ä¸ª|è¯¥|æ­¤)', '', source_name)
            source_name = re.sub(r'(?:æ–‡ä»¶å¤¹|æ–‡ä»¶|ç›®å½•)$', '', source_name).strip()
            
            # æ£€æŸ¥æ˜¯å¦åŒ…å«å¤šä¸ªå…³é”®è¯ï¼ˆç©ºæ ¼åˆ†éš”ï¼‰
            keywords = [kw.strip() for kw in source_name.split() if kw.strip()]
            
            logger.info(f"æŸ¥æ‰¾æºèŠ‚ç‚¹: åŸå§‹='{source_path}' â†’ æ¸…ç†å='{source_name}' â†’ å…³é”®è¯={keywords}")
            
            # å¦‚æœæœ‰å¤šä¸ªå…³é”®è¯ï¼Œè¿›è¡Œç»„åˆæŸ¥è¯¢
            if len(keywords) > 1:
                # å¤šå…³é”®è¯æŸ¥è¯¢ï¼šæ‰€æœ‰å…³é”®è¯éƒ½è¦åŒ¹é…
                conditions = []
                for keyword in keywords:
                    conditions.append(DocumentNode.name.like(f'%{keyword}%'))
                
                # æŸ¥æ‰¾åŒ…å«æ‰€æœ‰å…³é”®è¯çš„èŠ‚ç‚¹
                nodes = DocumentNode.query.filter(
                    and_(and_(*conditions), DocumentNode.is_deleted == False)
                ).all()
                
                if nodes:
                    logger.info(f"å¤šå…³é”®è¯æŸ¥è¯¢æ‰¾åˆ° {len(nodes)} ä¸ªåŒ¹é…èŠ‚ç‚¹")
                    # ä¼˜å…ˆè¿”å›æ–‡ä»¶å¤¹
                    for node in nodes:
                        if node.type == 'folder':
                            return node
                    return nodes[0]
                
                # å¦‚æœæ‰€æœ‰å…³é”®è¯éƒ½åŒ¹é…å¤±è´¥ï¼Œå°è¯•ä»»æ„å…³é”®è¯åŒ¹é…
                or_conditions = []
                for keyword in keywords:
                    or_conditions.append(DocumentNode.name.like(f'%{keyword}%'))
                
                nodes = DocumentNode.query.filter(
                    and_(or_(*or_conditions), DocumentNode.is_deleted == False)
                ).all()
                
                if nodes:
                    logger.info(f"ä»»æ„å…³é”®è¯æŸ¥è¯¢æ‰¾åˆ° {len(nodes)} ä¸ªåŒ¹é…èŠ‚ç‚¹")
                    # ä¼˜å…ˆè¿”å›æ–‡ä»¶å¤¹
                    for node in nodes:
                        if node.type == 'folder':
                            return node
                    return nodes[0]
            
            else:
                # å•å…³é”®è¯æŸ¥è¯¢ï¼šæŒ‰ä¼˜å…ˆçº§æŸ¥æ‰¾
                search_conditions = [
                    DocumentNode.name == source_name,  # å®Œå…¨åŒ¹é…
                    DocumentNode.name.like(f'%{source_name}%'),  # åŒ…å«åŒ¹é…
                ]
                
                for condition in search_conditions:
                    nodes = DocumentNode.query.filter(
                        and_(condition, DocumentNode.is_deleted == False)
                    ).all()
                    
                    if nodes:
                        logger.info(f"å•å…³é”®è¯æŸ¥è¯¢æ‰¾åˆ° {len(nodes)} ä¸ªåŒ¹é…èŠ‚ç‚¹")
                        # å¦‚æœæœ‰å¤šä¸ªåŒ¹é…ï¼Œä¼˜å…ˆè¿”å›æ–‡ä»¶å¤¹
                        for node in nodes:
                            if node.type == 'folder':
                                return node
                        # å¦‚æœæ²¡æœ‰æ–‡ä»¶å¤¹ï¼Œè¿”å›ç¬¬ä¸€ä¸ªæ–‡ä»¶
                        return nodes[0]
            
            return None
            
        except Exception as e:
            logger.error(f"æŸ¥æ‰¾æºèŠ‚ç‚¹å¤±è´¥: {e}")
            return None
    
    @staticmethod
    def _generate_from_folder(folder: DocumentNode, output_format: str, 
                             document_type: str, query_text: str, 
                             llm_model: str = None) -> Dict[str, Any]:
        """åŸºäºæ–‡ä»¶å¤¹ç”Ÿæˆæ–‡æ¡£"""
        try:
            logger.info(f"å¼€å§‹åŸºäºæ–‡ä»¶å¤¹ '{folder.name}' ç”Ÿæˆ {document_type} æ–‡æ¡£")
            
            # 1. è·å–æ–‡ä»¶å¤¹ä¸‹çš„æ‰€æœ‰æ–‡ä»¶
            files = DocumentGenerationService._get_folder_files(folder)
            if not files:
                return {
                    'success': False,
                    'error': f'æ–‡ä»¶å¤¹ "{folder.name}" ä¸­æ²¡æœ‰æ‰¾åˆ°æ–‡ä»¶',
                    'query': query_text,
                    'folder_info': folder.to_dict()
                }
            
            # 2. æå–æ–‡ä»¶å†…å®¹
            file_contents = []
            processed_count = 0
            
            for file_node in files:
                try:
                    content = DocumentGenerationService._extract_file_content(file_node)
                    if content:
                        file_contents.append({
                            'file_name': file_node.name,
                            'file_type': file_node.file_type,
                            'document_id': file_node.id,  # æ·»åŠ document_idç”¨äºç”Ÿæˆè¶…é“¾æ¥
                            'content': content[:2000]  # é™åˆ¶å†…å®¹é•¿åº¦é¿å…tokenè¿‡å¤š
                        })
                        processed_count += 1
                        
                        # æš‚æ—¶é™åˆ¶å¤„ç†æ–‡ä»¶æ•°é‡ï¼Œé¿å…tokenè¶…é™
                        if processed_count >= 10:
                            break
                except Exception as e:
                    logger.warning(f"æå–æ–‡ä»¶ {file_node.name} å†…å®¹å¤±è´¥: {e}")
                    continue
            
            if not file_contents:
                return {
                    'success': False,
                    'error': f'æ–‡ä»¶å¤¹ "{folder.name}" ä¸­çš„æ–‡ä»¶å†…å®¹æ— æ³•æå–',
                    'query': query_text,
                    'folder_info': folder.to_dict(),
                    'total_files': len(files)
                }
            
            # 3. ä½¿ç”¨LLMç”Ÿæˆæ–‡æ¡£
            generated_content = DocumentGenerationService._generate_content_with_llm(
                source_type='folder',
                source_name=folder.name,
                source_contents=file_contents,
                document_type=document_type,
                query_text=query_text,
                llm_model=llm_model
            )
            
            # 4. ä¿å­˜ç”Ÿæˆçš„æ–‡æ¡£ï¼ˆå½“å‰åªæ”¯æŒtxtï¼‰
            saved_file = DocumentGenerationService._save_generated_document(
                content=generated_content,
                source_name=folder.name,
                document_type=document_type,
                output_format=output_format,
                parent_folder=folder,
                query_text=query_text
            )
            
            return {
                'success': True,
                'query': query_text,
                'source_type': 'folder',
                'source_info': folder.to_dict(),
                'processed_files_count': processed_count,
                'total_files_count': len(files),
                'document_type': document_type,
                'output_format': output_format,
                'generated_content': generated_content,
                'saved_file': saved_file.to_dict() if saved_file else None,
                'generation_method': 'llm_folder_analysis',
                'timestamp': datetime.now().isoformat()
            }
            
        except Exception as e:
            logger.error(f"åŸºäºæ–‡ä»¶å¤¹ç”Ÿæˆæ–‡æ¡£å¤±è´¥: {e}")
            return {
                'success': False,
                'error': f'åŸºäºæ–‡ä»¶å¤¹ç”Ÿæˆæ–‡æ¡£å¤±è´¥: {str(e)}',
                'query': query_text,
                'folder_info': folder.to_dict() if folder else None
            }
    
    @staticmethod
    def _generate_from_file(file: DocumentNode, output_format: str, 
                           document_type: str, query_text: str,
                           llm_model: str = None) -> Dict[str, Any]:
        """åŸºäºå•ä¸ªæ–‡ä»¶ç”Ÿæˆæ–‡æ¡£"""
        try:
            logger.info(f"å¼€å§‹åŸºäºæ–‡ä»¶ '{file.name}' ç”Ÿæˆ {document_type} æ–‡æ¡£")
            
            # 1. æå–æ–‡ä»¶å†…å®¹
            content = DocumentGenerationService._extract_file_content(file)
            if not content:
                return {
                    'success': False,
                    'error': f'æ— æ³•æå–æ–‡ä»¶ "{file.name}" çš„å†…å®¹',
                    'query': query_text,
                    'file_info': file.to_dict()
                }
            
            # 2. å‡†å¤‡æ–‡ä»¶å†…å®¹æ•°æ®
            file_contents = [{
                'file_name': file.name,
                'file_type': file.file_type,
                'document_id': file.id,  # æ·»åŠ document_idç”¨äºç”Ÿæˆè¶…é“¾æ¥
                'content': content[:3000]  # å•æ–‡ä»¶å¯ä»¥å…è®¸æ›´å¤šå†…å®¹
            }]
            
            # 3. ä½¿ç”¨LLMç”Ÿæˆæ–‡æ¡£
            generated_content = DocumentGenerationService._generate_content_with_llm(
                source_type='file',
                source_name=file.name,
                source_contents=file_contents,
                document_type=document_type,
                query_text=query_text,
                llm_model=llm_model
            )
            
            # 4. ä¿å­˜ç”Ÿæˆçš„æ–‡æ¡£
            saved_file = DocumentGenerationService._save_generated_document(
                content=generated_content,
                source_name=file.name,
                document_type=document_type,
                output_format=output_format,
                parent_folder=file.parent if file.parent else None,
                query_text=query_text
            )
            
            return {
                'success': True,
                'query': query_text,
                'source_type': 'file',
                'source_info': file.to_dict(),
                'document_type': document_type,
                'output_format': output_format,
                'generated_content': generated_content,
                'saved_file': saved_file.to_dict() if saved_file else None,
                'generation_method': 'llm_file_analysis',
                'timestamp': datetime.now().isoformat()
            }
            
        except Exception as e:
            logger.error(f"åŸºäºæ–‡ä»¶ç”Ÿæˆæ–‡æ¡£å¤±è´¥: {e}")
            return {
                'success': False,
                'error': f'åŸºäºæ–‡ä»¶ç”Ÿæˆæ–‡æ¡£å¤±è´¥: {str(e)}',
                'query': query_text,
                'file_info': file.to_dict() if file else None
            }
    
    @staticmethod
    def _get_folder_files(folder: DocumentNode, max_depth: int = 2) -> List[DocumentNode]:
        """è·å–æ–‡ä»¶å¤¹ä¸‹çš„æ‰€æœ‰æ–‡ä»¶ï¼ˆé€’å½’ï¼‰"""
        try:
            files = []
            
            def collect_files(node: DocumentNode, current_depth: int = 0):
                if current_depth >= max_depth:
                    return
                
                for child in node.children:
                    if child.is_deleted:
                        continue
                        
                    if child.type == 'file':
                        files.append(child)
                    elif child.type == 'folder':
                        collect_files(child, current_depth + 1)
            
            collect_files(folder)
            return files
            
        except Exception as e:
            logger.error(f"è·å–æ–‡ä»¶å¤¹æ–‡ä»¶å¤±è´¥: {e}")
            return []
    
    @staticmethod
    def _extract_file_content(file: DocumentNode) -> Optional[str]:
        """æå–æ–‡ä»¶å†…å®¹ - æ”¯æŒå¤šç§æ–‡ä»¶æ ¼å¼"""
        try:
            # ä¼˜å…ˆä»DocumentContentè¡¨è·å–å·²æå–çš„å†…å®¹
            content_record = DocumentContent.query.filter_by(
                document_id=file.id
            ).first()
            
            if content_record and content_record.content_text:
                return content_record.content_text
            
            # æ£€æŸ¥æ–‡ä»¶è·¯å¾„æ˜¯å¦å­˜åœ¨
            if not file.file_path or not os.path.exists(file.file_path):
                logger.warning(f"æ–‡ä»¶è·¯å¾„ä¸å­˜åœ¨: {file.file_path}")
                return None
            
            file_type = file.file_type.lower() if file.file_type else ''
            
            try:
                # 1. PDFæ–‡ä»¶å¤„ç†
                if file_type in ['pdf']:
                    return DocumentGenerationService._extract_pdf_content(file.file_path)
                
                # 2. Wordæ–‡æ¡£å¤„ç†
                elif file_type in ['doc', 'docx']:
                    return DocumentGenerationService._extract_word_content(file.file_path)
                
                # 3. Excelæ–‡ä»¶å¤„ç†
                elif file_type in ['xls', 'xlsx']:
                    return DocumentGenerationService._extract_excel_content(file.file_path)
                
                # 4. PowerPointæ–‡ä»¶å¤„ç†
                elif file_type in ['ppt', 'pptx']:
                    return DocumentGenerationService._extract_ppt_content(file.file_path)
                
                # 5. å›¾ç‰‡æ–‡ä»¶OCRè¯†åˆ«
                elif file_type in ['jpg', 'jpeg', 'png', 'bmp', 'tiff', 'gif']:
                    ocr_result = ImageRecognitionService.recognize_text(file.file_path)
                    if ocr_result and ocr_result.get('success'):
                        return ocr_result.get('text', '')
                    return None
                
                # 6. æ–‡æœ¬æ–‡ä»¶ç›´æ¥è¯»å–
                elif file_type in ['txt', 'md', 'csv', 'json', 'xml', 'html', 'css', 'js', 'py', 'java', 'cpp', 'h']:
                    return DocumentGenerationService._extract_text_content(file.file_path)
                
                # 7. å…¶ä»–æ–‡æœ¬ç±»å‹æ–‡ä»¶å°è¯•è¯»å–
                else:
                    # å¯¹äºæœªçŸ¥ç±»å‹ï¼Œå°è¯•å½“ä½œæ–‡æœ¬æ–‡ä»¶è¯»å–
                    try:
                        return DocumentGenerationService._extract_text_content(file.file_path)
                    except:
                        logger.warning(f"æ— æ³•è¯†åˆ«æ–‡ä»¶ç±»å‹: {file.name} ({file_type})")
                        return None
                        
            except Exception as extract_error:
                logger.warning(f"æå–æ–‡ä»¶å†…å®¹å¤±è´¥ {file.name}: {extract_error}")
                return None
            
        except Exception as e:
            logger.error(f"æå–æ–‡ä»¶å†…å®¹å¤±è´¥ {file.name}: {e}")
            return None
    
    @staticmethod
    def _generate_content_with_llm(source_type: str, source_name: str, 
                                  source_contents: List[Dict], document_type: str,
                                  query_text: str, llm_model: str = None) -> str:
        """ä½¿ç”¨LLMç”Ÿæˆæ–‡æ¡£å†…å®¹"""
        try:
            # æ„å»ºç³»ç»Ÿæç¤ºè¯
            system_prompt = f"""ä½ æ˜¯ä¸€ä¸ªä¸“ä¸šçš„æ–‡æ¡£æ•´ç†ä¸“å®¶ã€‚ä½ çš„ä»»åŠ¡æ˜¯åŸºäºæä¾›çš„{source_type}å†…å®¹ï¼Œç”Ÿæˆä¸€ä»½{DocumentGenerationService.DOCUMENT_TYPES.get(document_type, 'é€šç”¨')}ç±»å‹çš„æ–‡æ¡£ã€‚

è¦æ±‚ï¼š
1. ä»”ç»†åˆ†ææ‰€æœ‰æä¾›çš„å†…å®¹ï¼Œæå–å…³é”®ä¿¡æ¯
2. æŒ‰ç…§{document_type}ç±»å‹çš„è¦æ±‚ç»„ç»‡å†…å®¹ç»“æ„
3. è¯­è¨€è¦ä¸“ä¸šã€å‡†ç¡®ã€ç®€æ´
4. ä¿æŒé€»è¾‘æ¸…æ™°ï¼Œæ¡ç†åˆ†æ˜
5. è¾“å‡ºçº¯æ–‡æœ¬æ ¼å¼ï¼Œä¸è¦ä½¿ç”¨markdownæˆ–å…¶ä»–æ ¼å¼æ ‡è®°
6. **é‡è¦ï¼šå½“æåˆ°å…·ä½“æ–‡ä»¶æ—¶ï¼Œè¯·ä½¿ç”¨[LINK:æ–‡ä»¶å:document_id]æ ¼å¼æ ‡è®°ï¼Œè¿™æ ·å¯ä»¥ä¸ºç”¨æˆ·æä¾›å¯ç‚¹å‡»çš„æ–‡ä»¶é“¾æ¥**

è¯·ç›´æ¥è¾“å‡ºç”Ÿæˆçš„æ–‡æ¡£å†…å®¹ï¼Œä¸è¦æ·»åŠ ä»»ä½•è§£é‡Šæˆ–è¯´æ˜ã€‚"""

            # æ„å»ºç”¨æˆ·æç¤ºè¯
            user_prompt = f"""è¯·åŸºäºä»¥ä¸‹{source_type} "{source_name}" çš„å†…å®¹ï¼Œç”Ÿæˆä¸€ä»½{DocumentGenerationService.DOCUMENT_TYPES.get(document_type, 'é€šç”¨')}æ–‡æ¡£ï¼š

ç”¨æˆ·éœ€æ±‚ï¼š{query_text}

{source_type}å†…å®¹ï¼š
"""
            
            # æ·»åŠ æ–‡ä»¶å†…å®¹
            for i, file_content in enumerate(source_contents, 1):
                user_prompt += f"""
æ–‡ä»¶{i}: {file_content['file_name']} ({file_content['file_type']}) [ID: {file_content.get('document_id', 'unknown')}]
å†…å®¹æ‘˜è¦:
{file_content['content'][:1500]}
{'...(å†…å®¹è¿‡é•¿å·²æˆªæ–­)' if len(file_content['content']) > 1500 else ''}

"""
            
            user_prompt += f"\nè¯·åŸºäºä»¥ä¸Šå†…å®¹ç”Ÿæˆ{DocumentGenerationService.DOCUMENT_TYPES.get(document_type, 'é€šç”¨')}æ–‡æ¡£ï¼š"
            
            # è°ƒç”¨LLMç”Ÿæˆå†…å®¹
            if not llm_model:
                from config import Config
                llm_model = f"{Config.DEFAULT_LLM_PROVIDER}:{Config.DEFAULT_LLM_MODEL}"
            
            logger.info(f"è°ƒç”¨LLMç”Ÿæˆæ–‡æ¡£ï¼Œæ¨¡å‹: {llm_model}")
            logger.debug(f"ç³»ç»Ÿæç¤ºè¯é•¿åº¦: {len(system_prompt)} å­—ç¬¦")
            logger.debug(f"ç”¨æˆ·æç¤ºè¯é•¿åº¦: {len(user_prompt)} å­—ç¬¦")
            
            # å°†ç³»ç»Ÿæç¤ºè¯å’Œç”¨æˆ·æç¤ºè¯åˆå¹¶
            combined_prompt = f"{system_prompt}\n\n{user_prompt}"
            
            response = LLMService.generate_answer(
                query=combined_prompt,
                context="",  # ä¸Šä¸‹æ–‡å·²åœ¨promptä¸­æä¾›
                llm_model=llm_model,
                scenario="document_generation"
            )
            
            if response and response.strip():
                logger.info(f"LLMç”ŸæˆæˆåŠŸï¼Œå†…å®¹é•¿åº¦: {len(response)} å­—ç¬¦")
                return response.strip()
            else:
                logger.warning("LLMè¿”å›ç©ºå“åº”ï¼Œè§¦å‘é™çº§å¤„ç†")
                return DocumentGenerationService._generate_simple_summary(source_contents, document_type)
            
        except Exception as e:
            logger.error(f"LLMç”Ÿæˆå†…å®¹å¤±è´¥: {e}")
            logger.error(f"é”™è¯¯è¯¦æƒ…: {str(e)}")
            # é™çº§å¤„ç†ï¼šç”Ÿæˆç®€å•çš„æ±‡æ€»
            logger.info("è§¦å‘é™çº§å¤„ç†ï¼šç”Ÿæˆç®€åŒ–æ±‡æ€»")
            return DocumentGenerationService._generate_simple_summary(source_contents, document_type)
    
    @staticmethod
    def _generate_simple_summary(source_contents: List[Dict], document_type: str) -> str:
        """ç”Ÿæˆç®€å•æ±‡æ€»ï¼ˆLLMå¤±è´¥æ—¶çš„é™çº§æ–¹æ¡ˆï¼‰"""
        try:
            logger.info(f"ç”Ÿæˆç®€åŒ–æ±‡æ€»ï¼Œæ–‡æ¡£ç±»å‹: {document_type}, æ–‡ä»¶æ•°é‡: {len(source_contents)}")
            
            summary_lines = [
                f"# {DocumentGenerationService.DOCUMENT_TYPES.get(document_type, 'æ–‡æ¡£æ±‡æ€»')}",
                f"ç”Ÿæˆæ—¶é—´ï¼š{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
                "",
                "## å†…å®¹æ¦‚è¿°",
                f"æœ¬æ–‡æ¡£åŸºäº {len(source_contents)} ä¸ªæ–‡ä»¶ç”Ÿæˆï¼ŒåŒ…å«ä»¥ä¸‹å†…å®¹ï¼š",
                ""
            ]
            
            # ç»Ÿè®¡æ–‡ä»¶ç±»å‹
            file_types = {}
            total_chars = 0
            
            for file_content in source_contents:
                file_type = file_content.get('file_type', 'æœªçŸ¥')
                file_types[file_type] = file_types.get(file_type, 0) + 1
                total_chars += len(file_content.get('content', ''))
            
            # æ·»åŠ ç»Ÿè®¡ä¿¡æ¯
            summary_lines.append("### æ–‡ä»¶ç»Ÿè®¡")
            for file_type, count in file_types.items():
                summary_lines.append(f"- {file_type}æ–‡ä»¶: {count} ä¸ª")
            summary_lines.append(f"- æ€»å†…å®¹é•¿åº¦: {total_chars:,} å­—ç¬¦")
            summary_lines.append("")
            
            # è¯¦ç»†æ–‡ä»¶åˆ—è¡¨
            summary_lines.append("## æ–‡ä»¶è¯¦æƒ…")
            
            for i, file_content in enumerate(source_contents, 1):
                file_name = file_content.get('file_name', 'æœªçŸ¥æ–‡ä»¶')
                file_type = file_content.get('file_type', 'æœªçŸ¥')
                document_id = file_content.get('document_id', None)
                content = file_content.get('content', '')
                
                # ä¸ºæ–‡ä»¶åæ·»åŠ è¶…é“¾æ¥æ ‡è®°ï¼ˆå¦‚æœæœ‰document_idï¼‰
                if document_id:
                    file_display_name = f"[LINK:{file_name}:{document_id}]"
                else:
                    file_display_name = file_name
                
                summary_lines.append(f"### {i}. {file_display_name}")
                summary_lines.append(f"**æ–‡ä»¶ç±»å‹:** {file_type}")
                summary_lines.append(f"**å†…å®¹é•¿åº¦:** {len(content)} å­—ç¬¦")
                
                # æå–å…³é”®å†…å®¹ç‰‡æ®µ
                if content:
                    # å–å¼€å¤´éƒ¨åˆ†
                    content_preview = content[:200].replace('\n', ' ').strip()
                    if content_preview:
                        summary_lines.append(f"**å†…å®¹é¢„è§ˆ:** {content_preview}...")
                    
                    # å¦‚æœå†…å®¹è¾ƒé•¿ï¼Œä¹Ÿæ˜¾ç¤ºä¸€äº›ä¸­é—´éƒ¨åˆ†
                    if len(content) > 400:
                        middle_content = content[len(content)//2:len(content)//2+150].replace('\n', ' ').strip()
                        if middle_content:
                            summary_lines.append(f"**ä¸­é—´ç‰‡æ®µ:** ...{middle_content}...")
                else:
                    summary_lines.append("**å†…å®¹é¢„è§ˆ:** (æ–‡ä»¶å†…å®¹ä¸ºç©º)")
                
                summary_lines.append("")
            
            # æ·»åŠ å»ºè®®æ€§è¯´æ˜
            summary_lines.extend([
                "## ä½¿ç”¨å»ºè®®",
                "æœ¬æ–‡æ¡£ä¸ºåŸºç¡€æ±‡æ€»ç‰ˆæœ¬ï¼ŒåŒ…å«äº†æºæ–‡ä»¶çš„ä¸»è¦ä¿¡æ¯å’Œå†…å®¹é¢„è§ˆã€‚",
                "æ‚¨å¯ä»¥æ ¹æ®ä»¥ä¸Šä¿¡æ¯äº†è§£æ–‡æ¡£çš„åŸºæœ¬ç»“æ„å’Œå†…å®¹æ¦‚å†µã€‚",
                "",
                "---",
                f"*æ–‡æ¡£ç”Ÿæˆäº {datetime.now().strftime('%Yå¹´%mæœˆ%dæ—¥ %H:%M:%S')}*"
            ])
            
            result = '\n'.join(summary_lines)
            logger.info(f"ç®€åŒ–æ±‡æ€»ç”Ÿæˆå®Œæˆï¼Œé•¿åº¦: {len(result)} å­—ç¬¦")
            return result
            
        except Exception as e:
            logger.error(f"ç”Ÿæˆç®€å•æ±‡æ€»å¤±è´¥: {e}")
            return f"""# æ–‡æ¡£ç”Ÿæˆå¤±è´¥

ç”Ÿæˆæ—¶é—´ï¼š{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}

## é”™è¯¯ä¿¡æ¯
æ–‡æ¡£ç”Ÿæˆè¿‡ç¨‹ä¸­é‡åˆ°é—®é¢˜ï¼š{str(e)}

## å»ºè®®
è¯·æ£€æŸ¥æºæ–‡ä»¶æ˜¯å¦å­˜åœ¨ï¼Œæˆ–è”ç³»ç³»ç»Ÿç®¡ç†å‘˜ã€‚"""
    
    @staticmethod
    def _save_generated_document(content: str, source_name: str, document_type: str,
                               output_format: str, parent_folder: DocumentNode = None, 
                               query_text: str = '') -> Optional[DocumentNode]:
        """ä¿å­˜ç”Ÿæˆçš„æ–‡æ¡£ - æ”¯æŒå¤šç§æ ¼å¼"""
        try:
            # ç”Ÿæˆè¯­ä¹‰åŒ–æ–‡ä»¶å
            filename = DocumentGenerationService._generate_semantic_filename(
                source_name, document_type, output_format, query_text)
            
            # ä¸ºç”Ÿæˆçš„æ–‡æ¡£åˆ›å»ºè™šæ‹Ÿè·¯å¾„ï¼ˆç”¨äºé¢„è§ˆç³»ç»Ÿè¯†åˆ«ï¼‰
            import uuid
            virtual_path = f"mcp_created/{uuid.uuid4()}_{filename}"
            
            # æ ¹æ®è¾“å‡ºæ ¼å¼ç”Ÿæˆå®é™…æ–‡ä»¶å†…å®¹å’ŒMIMEç±»å‹
            file_content, mime_type, file_size = DocumentGenerationService._generate_format_content(
                content, output_format)
            
            if file_content is None:
                logger.error(f"ç”Ÿæˆ {output_format} æ ¼å¼æ–‡ä»¶å¤±è´¥")
                return None
            
            # åˆ›å»ºDocumentNodeè®°å½•
            new_doc = DocumentNode(
                name=filename,
                type='file',
                parent_id=parent_folder.id if parent_folder else None,
                file_type=output_format,
                file_size=file_size,
                file_path=virtual_path,  # è®¾ç½®è™šæ‹Ÿè·¯å¾„ç”¨äºé¢„è§ˆ
                mime_type=mime_type,
                description=f"åŸºäº {source_name} ç”Ÿæˆçš„{DocumentGenerationService.DOCUMENT_TYPES.get(document_type, 'æ–‡æ¡£')}",
                created_by='document_generation_service',
                doc_metadata={
                    'generated_from': source_name,
                    'generation_type': document_type,
                    'generation_time': datetime.now().isoformat(),
                    'output_format': output_format
                }
            )
            
            db.session.add(new_doc)
            db.session.flush()  # è·å–ID
            
            # åˆ›å»ºDocumentContentè®°å½•
            doc_content = DocumentContent(
                document_id=new_doc.id,
                content_text=content,  # å§‹ç»ˆä¿å­˜åŸå§‹æ–‡æœ¬å†…å®¹ç”¨äºæœç´¢
                page_number=1,
                chunk_index=0,
                chunk_text=content[:1000]  # ä¿å­˜å‰1000ä¸ªå­—ç¬¦ä½œä¸ºchunk
            )
            
            db.session.add(doc_content)
            
            # å¯¹äºéæ–‡æœ¬æ ¼å¼ï¼Œéœ€è¦å°†äºŒè¿›åˆ¶å†…å®¹å­˜å‚¨åˆ°è™šæ‹Ÿæ–‡ä»¶ç³»ç»Ÿä¸­
            if output_format != 'txt' and file_content:
                DocumentGenerationService._store_virtual_file(virtual_path, file_content)
            
            db.session.commit()
            
            logger.info(f"æˆåŠŸä¿å­˜ç”Ÿæˆçš„ {output_format} æ–‡æ¡£: {filename} (ID: {new_doc.id})")
            return new_doc
            
        except Exception as e:
            logger.error(f"ä¿å­˜ç”Ÿæˆæ–‡æ¡£å¤±è´¥: {e}")
            db.session.rollback()
            return None
    
    @staticmethod
    def _generate_format_content(content: str, output_format: str) -> tuple[Optional[bytes], str, int]:
        """æ ¹æ®æ ¼å¼ç”Ÿæˆæ–‡ä»¶å†…å®¹
        
        Returns:
            tuple: (file_content, mime_type, file_size)
        """
        try:
            if output_format == 'txt':
                # æ–‡æœ¬æ ¼å¼
                file_content = content.encode('utf-8')
                return file_content, 'text/plain', len(file_content)
            
            elif output_format == 'pdf':
                # PDFæ ¼å¼
                return DocumentGenerationService._generate_pdf_content(content)
            
            elif output_format == 'xlsx':
                # Excelæ ¼å¼
                return DocumentGenerationService._generate_xlsx_content(content)
            
            elif output_format == 'docx':
                # Wordæ ¼å¼
                return DocumentGenerationService._generate_docx_content(content)
            
            else:
                logger.error(f"ä¸æ”¯æŒçš„è¾“å‡ºæ ¼å¼: {output_format}")
                return None, '', 0
                
        except Exception as e:
            logger.error(f"ç”Ÿæˆ {output_format} æ ¼å¼å†…å®¹å¤±è´¥: {e}")
            return None, '', 0
    
    @staticmethod
    def _generate_pdf_content(content: str) -> tuple[Optional[bytes], str, int]:
        """ç”ŸæˆPDFæ ¼å¼å†…å®¹"""
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont
            import io
            import os
            
            # åˆ›å»ºå†…å­˜ç¼“å†²åŒº
            buffer = io.BytesIO()
            
            # åˆ›å»ºPDFæ–‡æ¡£
            doc = SimpleDocTemplate(buffer, pagesize=A4)
            
            # è·å–æ ·å¼
            styles = getSampleStyleSheet()
            
            # å°è¯•æ³¨å†Œä¸­æ–‡å­—ä½“
            try:
                # æŒ‰ä¼˜å…ˆçº§æ’åˆ—å­—ä½“è·¯å¾„ï¼ˆä¼˜å…ˆé€‰æ‹©Unicodeæ”¯æŒå¥½çš„å­—ä½“ï¼‰
                font_paths = [
                    # é«˜ä¼˜å…ˆçº§ï¼šä¸“é—¨çš„Unicodeå­—ä½“
                    '/Library/Fonts/Arial Unicode MS.ttf',  # Arial Unicodeï¼ˆæœ€ä½³Unicodeæ”¯æŒï¼‰
                    'C:/Windows/Fonts/arialuni.ttf',  # Windows Arial Unicode
                    'C:/Windows/Fonts/tahoma.ttf',  # Tahomaï¼ˆè¾ƒå¥½çš„Unicodeæ”¯æŒï¼‰
                    
                    # ä¸­ç­‰ä¼˜å…ˆçº§ï¼šç³»ç»Ÿä¸­æ–‡å­—ä½“TTFç‰ˆæœ¬
                    'C:/Windows/Fonts/msyh.ttf',  # å¾®è½¯é›…é»‘TTFç‰ˆæœ¬
                    'C:/Windows/Fonts/simhei.ttf',  # é»‘ä½“TTF
                    'C:/Windows/Fonts/simsun.ttf',  # å®‹ä½“TTF
                    
                    # ä½ä¼˜å…ˆçº§ï¼šé€šç”¨ç³»ç»Ÿå­—ä½“
                    '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',  # Linux DejaVu
                    '/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf',  # Linux Liberation
                    'C:/Windows/Fonts/arial.ttf',  # Windows Arial
                    'C:/Windows/Fonts/calibri.ttf',  # Windows Calibri
                    
                    # æœ€åå°è¯•ï¼šmacOSå­—ä½“ï¼ˆéœ€è¦ç‰¹æ®Šå¤„ç†ï¼‰
                    '/System/Library/Fonts/Helvetica.ttc',  # macOS Helvetica
                    '/System/Library/Fonts/Arial.ttf',  # macOS Arialï¼ˆå¦‚æœå­˜åœ¨ï¼‰
                ]
                
                chinese_font_registered = False
                for font_path in font_paths:
                    if os.path.exists(font_path):
                        try:
                            font_name = 'UnicodeFont'
                            
                            # åªå¤„ç†TTFæ–‡ä»¶ï¼Œé¿å…TTCå…¼å®¹æ€§é—®é¢˜
                            if font_path.endswith('.ttf'):
                                pdfmetrics.registerFont(TTFont(font_name, font_path))
                            elif font_path.endswith('.ttc'):
                                # å¯¹äºTTCæ–‡ä»¶ï¼Œå°è¯•å¤šä¸ªå­å­—ä½“ç´¢å¼•
                                for subfont_index in [0, 1, 2]:
                                    try:
                                        pdfmetrics.registerFont(TTFont(font_name, font_path, subfontIndex=subfont_index))
                                        break
                                    except Exception:
                                        continue
                                else:
                                    # å¦‚æœæ‰€æœ‰å­å­—ä½“éƒ½å¤±è´¥ï¼Œè·³è¿‡è¿™ä¸ªæ–‡ä»¶
                                    raise Exception("All subfont indices failed")
                            else:
                                continue
                            
                            # åˆ›å»ºæ”¯æŒUnicodeçš„æ ·å¼
                            normal_style = ParagraphStyle(
                                'UnicodeNormal',
                                fontName=font_name,
                                fontSize=12,
                                leading=16,
                                spaceAfter=12,
                                encoding='utf-8'
                            )
                            chinese_font_registered = True
                            logger.info(f"æˆåŠŸæ³¨å†ŒUnicodeå­—ä½“: {font_path}")
                            break
                            
                        except Exception as font_error:
                            logger.warning(f"æ³¨å†Œå­—ä½“å¤±è´¥ {font_path}: {font_error}")
                            continue
                
                if not chinese_font_registered:
                    # æœ€åçš„å¤‡ç”¨æ–¹æ¡ˆï¼šä½¿ç”¨ReportLabå†…ç½®å­—ä½“ï¼Œå¹¶è¿›è¡Œæ–‡æœ¬é¢„å¤„ç†
                    logger.warning("æœªæ‰¾åˆ°å¯ç”¨çš„Unicodeå­—ä½“ï¼Œä½¿ç”¨å†…ç½®å­—ä½“å¹¶è¿›è¡Œæ–‡æœ¬é¢„å¤„ç†")
                    normal_style = styles['Normal']
                    # è®¾ç½®æ ‡è®°ï¼Œåç»­éœ€è¦å¯¹æ–‡æœ¬è¿›è¡Œé¢„å¤„ç†
                    normal_style._use_fallback = True
                else:
                    normal_style._use_fallback = False
                    
            except Exception as e:
                logger.error(f"å­—ä½“é…ç½®å¤±è´¥: {e}")
                normal_style = styles['Normal']
                normal_style._use_fallback = True
            
            # åˆ†å‰²å†…å®¹ä¸ºæ®µè½
            story = []
            paragraphs = content.split('\n\n')
            
            for para_text in paragraphs:
                if para_text.strip():
                    # å¤„ç†ç‰¹æ®Šå­—ç¬¦
                    para_text = para_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                    
                    # å¦‚æœä½¿ç”¨å¤‡ç”¨å­—ä½“ï¼Œè¿›è¡Œæ–‡æœ¬é¢„å¤„ç†
                    if hasattr(normal_style, '_use_fallback') and normal_style._use_fallback:
                        # å¯¹äºä¸æ”¯æŒUnicodeçš„å­—ä½“ï¼Œå°†ä¸­æ–‡å­—ç¬¦è½¬æ¢ä¸ºå¯æ˜¾ç¤ºçš„å½¢å¼
                        processed_text = ''
                        for char in para_text:
                            if ord(char) > 127:  # éASCIIå­—ç¬¦
                                # ä¿æŒä¸­æ–‡å­—ç¬¦ï¼Œè®©ReportLabå°è¯•å¤„ç†
                                processed_text += char
                            else:
                                processed_text += char
                        para_text = processed_text
                    
                    try:
                        para = Paragraph(para_text, normal_style)
                        story.append(para)
                        story.append(Spacer(1, 12))
                    except Exception as para_error:
                        # å¦‚æœæ®µè½åˆ›å»ºå¤±è´¥ï¼Œåˆ›å»ºä¸€ä¸ªç®€å•çš„æ–‡æœ¬æ®µè½
                        logger.warning(f"æ®µè½åˆ›å»ºå¤±è´¥ï¼Œä½¿ç”¨ç®€åŒ–å¤„ç†: {para_error}")
                        # ä½¿ç”¨åŸºç¡€æ ·å¼é‡è¯•
                        simple_style = styles['Normal']
                        try:
                            para = Paragraph(para_text.encode('ascii', 'replace').decode('ascii'), simple_style)
                            story.append(para)
                            story.append(Spacer(1, 12))
                        except Exception:
                            # æœ€åçš„å¤‡ç”¨æ–¹æ¡ˆï¼šçº¯æ–‡æœ¬
                            logger.warning("ä½¿ç”¨çº¯æ–‡æœ¬å¤‡ç”¨æ–¹æ¡ˆ")
                            from reportlab.platypus import Preformatted
                            para = Preformatted(para_text, simple_style)
                            story.append(para)
                            story.append(Spacer(1, 12))
            
            # æ„å»ºPDF
            doc.build(story)
            
            # è·å–å†…å®¹
            pdf_content = buffer.getvalue()
            buffer.close()
            
            return pdf_content, 'application/pdf', len(pdf_content)
            
        except ImportError:
            logger.error("ç¼ºå°‘PDFç”Ÿæˆä¾èµ–ï¼špip install reportlab")
            return None, '', 0
        except Exception as e:
            logger.error(f"ç”ŸæˆPDFå†…å®¹å¤±è´¥: {e}")
            return None, '', 0
    
    @staticmethod
    def _generate_xlsx_content(content: str) -> tuple[Optional[bytes], str, int]:
        """ç”ŸæˆExcelæ ¼å¼å†…å®¹"""
        try:
            import openpyxl
            from openpyxl.styles import Font, Alignment
            import io
            
            # åˆ›å»ºå·¥ä½œç°¿
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = "ç”Ÿæˆçš„æ–‡æ¡£"
            
            # è®¾ç½®æ ‡é¢˜æ ·å¼
            title_font = Font(bold=True, size=14)
            normal_font = Font(size=12)
            
            # åˆ†å‰²å†…å®¹ä¸ºè¡Œ
            lines = content.split('\n')
            row = 1
            
            for line in lines:
                if line.strip():
                    # æ£€æŸ¥æ˜¯å¦æ˜¯æ ‡é¢˜ï¼ˆä»¥#å¼€å¤´æˆ–åŒ…å«ç‰¹å®šå…³é”®è¯ï¼‰
                    if line.startswith('#') or any(keyword in line for keyword in ['æŠ¥å‘Š', 'æ€»ç»“', 'åˆ†æ']):
                        ws.cell(row=row, column=1, value=line.strip('#').strip())
                        ws.cell(row=row, column=1).font = title_font
                    else:
                        ws.cell(row=row, column=1, value=line.strip())
                        ws.cell(row=row, column=1).font = normal_font
                    
                    # è®¾ç½®è‡ªåŠ¨æ¢è¡Œ
                    ws.cell(row=row, column=1).alignment = Alignment(wrap_text=True)
                    row += 1
            
            # è®¾ç½®åˆ—å®½
            ws.column_dimensions['A'].width = 80
            
            # ä¿å­˜åˆ°å†…å­˜
            buffer = io.BytesIO()
            wb.save(buffer)
            excel_content = buffer.getvalue()
            buffer.close()
            
            return excel_content, 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', len(excel_content)
            
        except ImportError:
            logger.error("ç¼ºå°‘Excelç”Ÿæˆä¾èµ–ï¼špip install openpyxl")
            return None, '', 0
        except Exception as e:
            logger.error(f"ç”ŸæˆExcelå†…å®¹å¤±è´¥: {e}")
            return None, '', 0
    
    @staticmethod
    def _generate_docx_content(content: str) -> tuple[Optional[bytes], str, int]:
        """ç”ŸæˆWordæ ¼å¼å†…å®¹"""
        try:
            from docx import Document
            from docx.shared import Inches
            import io
            
            # åˆ›å»ºæ–‡æ¡£
            doc = Document()
            
            # åˆ†å‰²å†…å®¹ä¸ºæ®µè½
            paragraphs = content.split('\n\n')
            
            for para_text in paragraphs:
                if para_text.strip():
                    # æ£€æŸ¥æ˜¯å¦æ˜¯æ ‡é¢˜
                    if para_text.startswith('#'):
                        # å¤„ç†ä¸åŒçº§åˆ«çš„æ ‡é¢˜
                        level = len(para_text) - len(para_text.lstrip('#'))
                        title_text = para_text.lstrip('#').strip()
                        if level == 1:
                            doc.add_heading(title_text, level=1)
                        elif level == 2:
                            doc.add_heading(title_text, level=2)
                        else:
                            doc.add_heading(title_text, level=3)
                    else:
                        # æ™®é€šæ®µè½
                        doc.add_paragraph(para_text.strip())
            
            # ä¿å­˜åˆ°å†…å­˜
            buffer = io.BytesIO()
            doc.save(buffer)
            docx_content = buffer.getvalue()
            buffer.close()
            
            return docx_content, 'application/vnd.openxmlformats-officedocument.wordprocessingml.document', len(docx_content)
            
        except ImportError:
            logger.error("ç¼ºå°‘Wordç”Ÿæˆä¾èµ–ï¼špip install python-docx")
            return None, '', 0
        except Exception as e:
            logger.error(f"ç”ŸæˆWordå†…å®¹å¤±è´¥: {e}")
            return None, '', 0
    
    @staticmethod
    def _generate_semantic_filename(source_name: str, document_type: str, 
                                   output_format: str, query_text: str = '') -> str:
        """ç”Ÿæˆè¯­ä¹‰åŒ–çš„æ–‡ä»¶å"""
        try:
            # æ¸…ç†æºæ–‡ä»¶åï¼ˆå»æ‰æ‰©å±•åï¼‰
            clean_source = re.sub(r'\.[a-zA-Z0-9]+$', '', source_name)
            
            # ä»æŸ¥è¯¢æ–‡æœ¬ä¸­æå–è¯­ä¹‰å…³é”®è¯
            semantic_suffix = DocumentGenerationService._extract_semantic_suffix(query_text, document_type)
            
            # æ„å»ºè¯­ä¹‰åŒ–æ–‡ä»¶å
            if semantic_suffix:
                filename = f"{clean_source}{semantic_suffix}.{output_format}"
            else:
                # ä½¿ç”¨é»˜è®¤çš„æ–‡æ¡£ç±»å‹åç§°
                type_name = DocumentGenerationService.DOCUMENT_TYPES.get(document_type, 'æ–‡æ¡£')
                filename = f"{clean_source}{type_name}.{output_format}"
            
            # ç¡®ä¿æ–‡ä»¶åä¸è¶…è¿‡255ä¸ªå­—ç¬¦
            if len(filename) > 255:
                max_source_len = 255 - len(semantic_suffix or type_name) - len(output_format) - 2
                clean_source = clean_source[:max_source_len]
                filename = f"{clean_source}{semantic_suffix or type_name}.{output_format}"
            
            return filename
            
        except Exception as e:
            logger.error(f"ç”Ÿæˆè¯­ä¹‰åŒ–æ–‡ä»¶åå¤±è´¥: {e}")
            # è¿”å›ç®€å•çš„é»˜è®¤åç§°
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            return f"{source_name}_{document_type}_{timestamp}.{output_format}"
    
    @staticmethod
    def _extract_semantic_suffix(query_text: str, document_type: str) -> str:
        """ä»æŸ¥è¯¢æ–‡æœ¬ä¸­æå–è¯­ä¹‰åŒ–åç¼€"""
        try:
            # å®šä¹‰è¯­ä¹‰å…³é”®è¯æ˜ å°„
            semantic_patterns = {
                # ä»·æ ¼ç›¸å…³
                r'ä»·æ ¼(?:é¢„æµ‹|åˆ†æ|è¯„ä¼°|ä¼°ç®—)': 'ä»·æ ¼é¢„æµ‹æŠ¥å‘Š',
                r'(?:æˆ¿ä»·|ä»·å€¼)(?:é¢„æµ‹|åˆ†æ|è¯„ä¼°)': 'ä»·æ ¼é¢„æµ‹æŠ¥å‘Š',
                r'å¸‚åœº(?:ä»·æ ¼|è¡Œæƒ…)(?:åˆ†æ|é¢„æµ‹)': 'å¸‚åœºè¡Œæƒ…åˆ†æ',
                
                # åˆ†æç›¸å…³
                r'(?:æ•°æ®|å†…å®¹|æƒ…å†µ)åˆ†æ': 'æ•°æ®åˆ†ææŠ¥å‘Š',
                r'(?:å¸‚åœº|è¡Œä¸š|ç«äº‰)åˆ†æ': 'å¸‚åœºåˆ†ææŠ¥å‘Š',
                r'(?:è´¢åŠ¡|ç»è¥|ä¸šåŠ¡)åˆ†æ': 'è´¢åŠ¡åˆ†ææŠ¥å‘Š',
                r'(?:é£é™©|å®‰å…¨)(?:åˆ†æ|è¯„ä¼°)': 'é£é™©è¯„ä¼°æŠ¥å‘Š',
                
                # æ€»ç»“ç›¸å…³
                r'(?:å·¥ä½œ|é¡¹ç›®|ä¸šåŠ¡)æ€»ç»“': 'å·¥ä½œæ€»ç»“',
                r'(?:ä¼šè®®|æ´»åŠ¨|äº‹ä»¶)æ€»ç»“': 'ä¼šè®®æ€»ç»“',
                r'(?:å­¦ä¹ |åŸ¹è®­)æ€»ç»“': 'å­¦ä¹ æ€»ç»“',
                
                # æŠ¥å‘Šç›¸å…³
                r'(?:è°ƒç ”|è°ƒæŸ¥)æŠ¥å‘Š': 'è°ƒç ”æŠ¥å‘Š',
                r'(?:æµ‹è¯•|è¯„æµ‹)æŠ¥å‘Š': 'æµ‹è¯•æŠ¥å‘Š',
                r'(?:è¿›å±•|è¿›åº¦)æŠ¥å‘Š': 'è¿›åº¦æŠ¥å‘Š',
                r'(?:å¹´åº¦|æœˆåº¦|å­£åº¦)æŠ¥å‘Š': 'å®šæœŸæŠ¥å‘Š',
                
                # è®¡åˆ’ç›¸å…³
                r'(?:å·¥ä½œ|é¡¹ç›®|è¥é”€)è®¡åˆ’': 'å·¥ä½œè®¡åˆ’',
                r'(?:å‘å±•|æˆ˜ç•¥)è§„åˆ’': 'å‘å±•è§„åˆ’',
                
                # æ–¹æ¡ˆç›¸å…³
                r'(?:è§£å†³|å®æ–½|ä¼˜åŒ–)æ–¹æ¡ˆ': 'è§£å†³æ–¹æ¡ˆ',
                r'(?:æ”¹è¿›|æå‡)æ–¹æ¡ˆ': 'æ”¹è¿›æ–¹æ¡ˆ'
            }
            
            query_lower = query_text.lower()
            
            # æŒ‰ä¼˜å…ˆçº§æ£€æŸ¥è¯­ä¹‰æ¨¡å¼
            for pattern, suffix in semantic_patterns.items():
                if re.search(pattern, query_text):
                    return suffix
            
            # å¦‚æœæ²¡æœ‰æ‰¾åˆ°ç‰¹å®šæ¨¡å¼ï¼Œå°è¯•æå–åŠ¨è¯+åè¯ç»„åˆ
            action_patterns = {
                r'(é¢„æµ‹|åˆ†æ|è¯„ä¼°|ä¼°ç®—)(.{1,8}?)(?:æŠ¥å‘Š|æ–‡æ¡£|èµ„æ–™)': r'\1\2æŠ¥å‘Š',
                r'(æ€»ç»“|æ±‡æ€»|æ¢³ç†)(.{1,8}?)(?:æ–‡æ¡£|èµ„æ–™|å†…å®¹)': r'\1\2',
                r'(åˆ¶ä½œ|ç”Ÿæˆ|ç¼–å†™)(.{1,15}?)(?:æŠ¥å‘Š|æ–‡æ¡£)': r'\2æŠ¥å‘Š'
            }
            
            for pattern, replacement in action_patterns.items():
                match = re.search(pattern, query_text)
                if match:
                    return re.sub(pattern, replacement, match.group(0))
            
            # é»˜è®¤æ ¹æ®æ–‡æ¡£ç±»å‹è¿”å›
            type_mappings = {
                'report': 'åˆ†ææŠ¥å‘Š',
                'summary': 'æ€»ç»“æŠ¥å‘Š', 
                'analysis': 'åˆ†ææ–‡æ¡£',
                'documentation': 'è¯´æ˜æ–‡æ¡£',
                'other': 'æ–‡æ¡£'
            }
            
            return type_mappings.get(document_type, 'æ–‡æ¡£')
            
        except Exception as e:
            logger.error(f"æå–è¯­ä¹‰åç¼€å¤±è´¥: {e}")
            return DocumentGenerationService.DOCUMENT_TYPES.get(document_type, 'æ–‡æ¡£')
    
    @staticmethod
    def _extract_text_content(file_path: str) -> Optional[str]:
        """æå–æ–‡æœ¬æ–‡ä»¶å†…å®¹"""
        try:
            # å°è¯•UTF-8ç¼–ç è¯»å–
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # å°è¯•å…¶ä»–ç¼–ç 
            for encoding in ['gbk', 'gb2312', 'latin1', 'iso-8859-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        return f.read()
                except:
                    continue
            logger.warning(f"æ— æ³•è¯†åˆ«æ–‡æœ¬æ–‡ä»¶ç¼–ç : {file_path}")
            return None
        except Exception as e:
            logger.error(f"è¯»å–æ–‡æœ¬æ–‡ä»¶å¤±è´¥ {file_path}: {e}")
            return None
    
    @staticmethod
    def _extract_pdf_content(file_path: str) -> Optional[str]:
        """æå–PDFæ–‡ä»¶å†…å®¹"""
        try:
            # ä½¿ç”¨é¡¹ç›®ä¸­çš„PDFå‘é‡åŒ–å™¨
            from app.services.vectorization.pdf_vectorizer import PDFVectorizer
            vectorizer = PDFVectorizer()
            result = vectorizer.extract_text(file_path, extract_images=False)  # ä¸æå–å›¾ç‰‡ä»¥åŠ å¿«é€Ÿåº¦
            if result and result.get('success'):
                return result.get('text', '')
            
            # å¦‚æœPDFå‘é‡åŒ–å™¨å¤±è´¥ï¼Œå°è¯•ä½¿ç”¨PyPDF2ä½œä¸ºé™çº§æ–¹æ¡ˆ
            try:
                import PyPDF2
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    text = ""
                    for page in pdf_reader.pages:
                        text += page.extract_text() + "\n"
                    return text.strip()
            except ImportError:
                logger.warning("PyPDF2æœªå®‰è£…ï¼Œè¯·å®‰è£…ï¼špip install PyPDF2")
            except Exception as pypdf_error:
                logger.warning(f"PyPDF2æå–å¤±è´¥: {pypdf_error}")
            
            # å°è¯•ä½¿ç”¨pdfplumber
            try:
                import pdfplumber
                with pdfplumber.open(file_path) as pdf:
                    text = ""
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                    return text.strip()
            except ImportError:
                logger.warning("pdfplumberæœªå®‰è£…ï¼Œè¯·å®‰è£…ï¼špip install pdfplumber")
            except Exception as plumber_error:
                logger.warning(f"pdfplumberæå–å¤±è´¥: {plumber_error}")
            
            return None
            
        except Exception as e:
            logger.error(f"æå–PDFå†…å®¹å¤±è´¥ {file_path}: {e}")
            return None
    
    @staticmethod
    def _extract_word_content(file_path: str) -> Optional[str]:
        """æå–Wordæ–‡æ¡£å†…å®¹"""
        try:
            # ä½¿ç”¨é¡¹ç›®ä¸­çš„Wordå‘é‡åŒ–å™¨
            from app.services.vectorization.word_vectorizer import WordVectorizer
            vectorizer = WordVectorizer()
            result = vectorizer.extract_text(file_path)
            if result and result.get('success'):
                return result.get('text', '')
            
            # å¦‚æœWordå‘é‡åŒ–å™¨å¤±è´¥ï¼Œå°è¯•ä½¿ç”¨python-docxä½œä¸ºé™çº§æ–¹æ¡ˆ
            try:
                import docx
                doc = docx.Document(file_path)
                text = ""
                for paragraph in doc.paragraphs:
                    text += paragraph.text + "\n"
                
                # æå–è¡¨æ ¼å†…å®¹
                for table in doc.tables:
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            row_text.append(cell.text.strip())
                        text += "\t".join(row_text) + "\n"
                
                return text.strip()
            except ImportError:
                logger.warning("python-docxæœªå®‰è£…ï¼Œè¯·å®‰è£…ï¼špip install python-docx")
            except Exception as docx_error:
                logger.warning(f"python-docxæå–å¤±è´¥: {docx_error}")
            
            return None
            
        except Exception as e:
            logger.error(f"æå–Wordå†…å®¹å¤±è´¥ {file_path}: {e}")
            return None
    
    @staticmethod
    def _extract_excel_content(file_path: str) -> Optional[str]:
        """æå–Excelæ–‡ä»¶å†…å®¹"""
        try:
            # ä½¿ç”¨é¡¹ç›®ä¸­çš„Excelå‘é‡åŒ–å™¨
            from app.services.vectorization.excel_vectorizer import ExcelVectorizer
            vectorizer = ExcelVectorizer()
            result = vectorizer.extract_text(file_path)
            if result and result.get('success'):
                return result.get('text', '')
            
            # å¦‚æœExcelå‘é‡åŒ–å™¨å¤±è´¥ï¼Œå°è¯•ä½¿ç”¨pandasä½œä¸ºé™çº§æ–¹æ¡ˆ
            try:
                import pandas as pd
                
                # è¯»å–æ‰€æœ‰å·¥ä½œè¡¨
                excel_file = pd.ExcelFile(file_path)
                text = ""
                
                for sheet_name in excel_file.sheet_names:
                    text += f"=== å·¥ä½œè¡¨: {sheet_name} ===\n"
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    
                    # è½¬æ¢ä¸ºæ–‡æœ¬æ ¼å¼
                    text += df.to_string(index=False, na_rep='') + "\n\n"
                
                return text.strip()
            except ImportError:
                logger.warning("pandasæˆ–openpyxlæœªå®‰è£…ï¼Œè¯·å®‰è£…ï¼špip install pandas openpyxl")
            except Exception as excel_error:
                logger.warning(f"pandasæå–Excelå¤±è´¥: {excel_error}")
            
            # å°è¯•ä½¿ç”¨xlrdï¼ˆå¯¹äºè€ç‰ˆæœ¬çš„.xlsæ–‡ä»¶ï¼‰
            try:
                import xlrd
                workbook = xlrd.open_workbook(file_path)
                text = ""
                
                for sheet_name in workbook.sheet_names():
                    sheet = workbook.sheet_by_name(sheet_name)
                    text += f"=== å·¥ä½œè¡¨: {sheet_name} ===\n"
                    
                    for row_idx in range(sheet.nrows):
                        row_values = []
                        for col_idx in range(sheet.ncols):
                            cell_value = sheet.cell_value(row_idx, col_idx)
                            row_values.append(str(cell_value))
                        text += "\t".join(row_values) + "\n"
                    text += "\n"
                
                return text.strip()
            except ImportError:
                logger.warning("xlrdæœªå®‰è£…ï¼Œè¯·å®‰è£…ï¼špip install xlrd")
            except Exception as xlrd_error:
                logger.warning(f"xlrdæå–å¤±è´¥: {xlrd_error}")
            
            return None
            
        except Exception as e:
            logger.error(f"æå–Excelå†…å®¹å¤±è´¥ {file_path}: {e}")
            return None
    
    @staticmethod
    def _extract_ppt_content(file_path: str) -> Optional[str]:
        """æå–PowerPointæ–‡ä»¶å†…å®¹"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_content = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
            return '\n'.join(text_content)
        except ImportError:
            logger.warning("ç¼ºå°‘PowerPointå¤„ç†åº“ï¼špip install python-pptx")
            return None
        except Exception as e:
            logger.error(f"æå–PowerPointæ–‡ä»¶å†…å®¹å¤±è´¥ {file_path}: {e}")
            return None
    
    @staticmethod
    def _store_virtual_file(virtual_path: str, file_content: bytes) -> bool:
        """å­˜å‚¨è™šæ‹Ÿæ–‡ä»¶å†…å®¹åˆ°å†…å­˜æˆ–ä¸´æ—¶å­˜å‚¨"""
        try:
            import os
            
            # ç¡®ä¿è™šæ‹Ÿæ–‡ä»¶ç›®å½•å­˜åœ¨
            virtual_dir = os.path.dirname(virtual_path)
            full_virtual_dir = os.path.join('uploads', virtual_dir)
            
            if not os.path.exists(full_virtual_dir):
                os.makedirs(full_virtual_dir, exist_ok=True)
            
            # å°†æ–‡ä»¶å†…å®¹å†™å…¥è™šæ‹Ÿè·¯å¾„
            full_virtual_path = os.path.join('uploads', virtual_path)
            with open(full_virtual_path, 'wb') as f:
                f.write(file_content)
            
            logger.info(f"è™šæ‹Ÿæ–‡ä»¶å·²å­˜å‚¨: {full_virtual_path}")
            return True
            
        except Exception as e:
            logger.error(f"å­˜å‚¨è™šæ‹Ÿæ–‡ä»¶å¤±è´¥ {virtual_path}: {e}")
            return False 